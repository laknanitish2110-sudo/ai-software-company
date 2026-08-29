from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import math
import random
import hashlib

PROJECTS_DIR = Path("generated_projects")
ASSETS_DIR = Path(__file__).parent.parent.parent / "assets"
TEMPLATE_DIR = ASSETS_DIR / "templates"

TEMPLATE_FILES = [
    "dark-modern.pptx",
    "dark-dynamic-lines.pptx",
    "hendrix-dark-gradient.pptx",
]

_use_template_bg = False
_blank_layout_idx = 6

# ─── Design System ───────────────────────────────────────────────
# Dark gradient palette inspired by Dribbble/Pinterest tech pitch decks

# Primary backgrounds
BG_DEEP = (8, 4, 24)        # Deep space purple
BG_DARK = (12, 8, 32)       # Dark purple
BG_SECTION = (16, 12, 40)   # Section purple
BG_CARD = (24, 18, 56)      # Card/elevated surface

# Gradients (top-left to bottom-right feel)
GRAD_TOP = "08041A"          # Near-black purple
GRAD_BOTTOM = "1A1040"       # Dark indigo
GRAD_ACCENT_START = "635BFF" # Primary accent
GRAD_ACCENT_END = "4F46E5"   # Deeper accent
GRAD_WARM_START = "7C3AED"   # Violet
GRAD_WARM_END = "4F46E5"     # Indigo

# Accent colors
ACCENT = (99, 91, 255)       # Primary purple #635BFF
ACCENT_LIGHT = (139, 131, 255)
ACCENT_DIM = (60, 55, 140)
TEAL = (11, 191, 140)        # Success/highlight #0BBF8C
CORAL = (237, 95, 116)       # Warning/emphasis #ED5F74
AMBER = (245, 166, 35)       # Stats/numbers #F5A623
CYAN = (0, 181, 216)         # Info/links #00B5D8

# Text colors
WHITE = (250, 250, 250)
TEXT_PRIMARY = (240, 240, 245)
TEXT_SECONDARY = (180, 180, 200)
TEXT_MUTED = (120, 120, 150)

# Typography
FONT_HEADING = "Poppins"
FONT_BODY = "Poppins"
FONT_MONO = "Consolas"

# Slide dimensions (16:9 widescreen)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Layout grid
MARGIN_L = Inches(0.9)
MARGIN_R = Inches(0.9)
MARGIN_T = Inches(0.8)
CONTENT_W = Inches(11.533)
CONTENT_H = Inches(5.5)


def _rgb(r, g, b):
    return RGBColor(r, g, b)


# ─── Template Helpers ───────────────────────────────────────────

def _find_blank_layout(prs):
    for i, layout in enumerate(prs.slide_layouts):
        name = layout.name.upper()
        if 'BLANK' in name and len(layout.placeholders) <= 1:
            return i
    for i, layout in enumerate(prs.slide_layouts):
        if len(layout.placeholders) == 0:
            return i
    return len(prs.slide_layouts) - 1


def _clear_all_slides(prs):
    prs_element = prs.part._element
    sldIdLst = prs_element.find(qn('p:sldIdLst'))
    if sldIdLst is not None:
        for sldId in list(sldIdLst):
            rId = sldId.get(qn('r:id'))
            prs.part.drop_rel(rId)
            sldIdLst.remove(sldId)


def _create_presentation(project_id: str = ""):
    global _use_template_bg, _blank_layout_idx

    templates = [TEMPLATE_DIR / tf for tf in TEMPLATE_FILES if (TEMPLATE_DIR / tf).exists()]

    seed = int(hashlib.md5(project_id.encode()).hexdigest()[:8], 16)
    choice = seed % (len(templates) + 1)

    if choice == 0 or not templates:
        prs = Presentation()
        _use_template_bg = False
        _blank_layout_idx = 6
    else:
        template_path = templates[choice - 1]
        prs = Presentation(str(template_path))
        _clear_all_slides(prs)
        _use_template_bg = True
        _blank_layout_idx = _find_blank_layout(prs)

    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    return prs


def _get_template_name(project_id: str = "") -> str:
    templates = [tf for tf in TEMPLATE_FILES if (TEMPLATE_DIR / tf).exists()]
    if not templates:
        return "custom"
    seed = int(hashlib.md5(project_id.encode()).hexdigest()[:8], 16)
    choice = seed % (len(templates) + 1)
    if choice == 0:
        return "custom"
    return templates[choice - 1].replace(".pptx", "")


# ─── Background Helpers ─────────────────────────────────────────

def _set_gradient_bg(slide, color_top, color_bottom, angle=5400000):
    bg = slide.background._element
    bg_ref = bg.find(qn('p:bgRef'))
    if bg_ref is not None:
        bg.remove(bg_ref)
    bg_pr = bg.find(qn('p:bgPr'))
    if bg_pr is not None:
        bg.remove(bg_pr)

    bg_pr = etree.SubElement(bg, qn('p:bgPr'))
    grad = etree.SubElement(bg_pr, qn('a:gradFill'))
    gs_lst = etree.SubElement(grad, qn('a:gsLst'))

    gs1 = etree.SubElement(gs_lst, qn('a:gs'))
    gs1.set('pos', '0')
    c1 = etree.SubElement(gs1, qn('a:srgbClr'))
    c1.set('val', color_top)

    gs2 = etree.SubElement(gs_lst, qn('a:gs'))
    gs2.set('pos', '100000')
    c2 = etree.SubElement(gs2, qn('a:srgbClr'))
    c2.set('val', color_bottom)

    lin = etree.SubElement(grad, qn('a:lin'))
    lin.set('ang', str(angle))
    lin.set('scaled', '1')

    etree.SubElement(bg_pr, qn('a:effectLst'))


def _set_three_stop_gradient_bg(slide, c1_hex, c2_hex, c3_hex, angle=5400000):
    bg = slide.background._element
    bg_ref = bg.find(qn('p:bgRef'))
    if bg_ref is not None:
        bg.remove(bg_ref)
    bg_pr = bg.find(qn('p:bgPr'))
    if bg_pr is not None:
        bg.remove(bg_pr)

    bg_pr = etree.SubElement(bg, qn('p:bgPr'))
    grad = etree.SubElement(bg_pr, qn('a:gradFill'))
    gs_lst = etree.SubElement(grad, qn('a:gsLst'))

    for pos, color in [(0, c1_hex), (50000, c2_hex), (100000, c3_hex)]:
        gs = etree.SubElement(gs_lst, qn('a:gs'))
        gs.set('pos', str(pos))
        c = etree.SubElement(gs, qn('a:srgbClr'))
        c.set('val', color)

    lin = etree.SubElement(grad, qn('a:lin'))
    lin.set('ang', str(angle))
    lin.set('scaled', '1')
    etree.SubElement(bg_pr, qn('a:effectLst'))


# ─── Decorative Shape Helpers ────────────────────────────────────

def _add_top_accent_bar(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(*ACCENT)
    shape.line.fill.background()
    shape.shadow.inherit = False


def _add_bottom_line(slide):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_L, SLIDE_H - Inches(0.6),
        CONTENT_W, Inches(0.01)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(40, 35, 80)
    shape.line.fill.background()
    shape.shadow.inherit = False


def _add_corner_decoration(slide, position="top-right"):
    if position == "top-right":
        x, y = SLIDE_W - Inches(2.5), Inches(0.3)
    elif position == "bottom-left":
        x, y = Inches(0.3), SLIDE_H - Inches(2.5)
    else:
        x, y = SLIDE_W - Inches(2.5), SLIDE_H - Inches(2.5)

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(2), Inches(2))
    fill = circle.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(*ACCENT)

    _set_shape_alpha(circle, 6000)
    circle.line.fill.background()
    circle.shadow.inherit = False


def _set_shape_alpha(shape, alpha_pct):
    sp_pr = shape._element.find(qn('p:spPr'))
    if sp_pr is None:
        for el in shape._element.iter():
            if el.tag.endswith('spPr'):
                sp_pr = el
                break
    if sp_pr is not None:
        solid = sp_pr.find(qn('a:solidFill'))
        if solid is not None:
            clr = solid.find(qn('a:srgbClr'))
            if clr is not None:
                for old in clr.findall(qn('a:alpha')):
                    clr.remove(old)
                alpha = etree.SubElement(clr, qn('a:alpha'))
                alpha.set('val', str(alpha_pct))


def _add_glow_circle(slide, x, y, size, color_rgb, alpha_pct=4000):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = _rgb(*color_rgb)
    _set_shape_alpha(circle, alpha_pct)
    circle.line.fill.background()
    circle.shadow.inherit = False


def _add_left_accent_stripe(slide, top=None, height=None):
    t = top or Inches(1.6)
    h = height or Inches(4.5)
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), t, Inches(0.05), h)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = _rgb(*ACCENT)
    stripe.line.fill.background()
    stripe.shadow.inherit = False


def _add_slide_number(slide, num, total):
    left = SLIDE_W - Inches(1.5)
    top = SLIDE_H - Inches(0.55)
    txBox = slide.shapes.add_textbox(left, top, Inches(1.2), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{num:02d}"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = _rgb(*ACCENT)
    run.font.name = FONT_MONO

    run2 = p.add_run()
    run2.text = f" / {total:02d}"
    run2.font.size = Pt(11)
    run2.font.color.rgb = _rgb(*TEXT_MUTED)
    run2.font.name = FONT_MONO
    p.alignment = PP_ALIGN.RIGHT


def _add_footer_branding(slide):
    txBox = slide.shapes.add_textbox(MARGIN_L, SLIDE_H - Inches(0.55), Inches(4), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "AI SOFTWARE COMPANY"
    run.font.size = Pt(8)
    run.font.bold = True
    run.font.color.rgb = _rgb(*TEXT_MUTED)
    run.font.name = FONT_HEADING
    run.font.spacing = Pt(3)


# ─── Text Helpers ────────────────────────────────────────────────

def _add_title_text(slide, text, left=None, top=None, width=None, size=36, color=TEXT_PRIMARY, bold=True):
    l = left or MARGIN_L
    t = top or Inches(0.6)
    w = width or CONTENT_W
    txBox = slide.shapes.add_textbox(l, t, w, Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = _rgb(*color)
    p.font.name = FONT_HEADING
    p.line_spacing = Pt(size * 1.2)
    return txBox


def _add_subtitle_text(slide, text, left=None, top=None, width=None, size=16, color=TEXT_SECONDARY):
    l = left or MARGIN_L
    t = top or Inches(1.8)
    w = width or CONTENT_W
    txBox = slide.shapes.add_textbox(l, t, w, Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(size)
    p.font.color.rgb = _rgb(*color)
    p.font.name = FONT_BODY
    return txBox


def _add_bullet_content(slide, items, left=None, top=None, width=None, height=None,
                        bullet_color=ACCENT, text_color=TEXT_SECONDARY, text_size=17, spacing=14):
    l = left or Inches(1.1)
    t = top or Inches(2.0)
    w = width or Inches(10.5)
    h = height or Inches(4.8)
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True

    items = items if isinstance(items, list) else [str(items)]
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        text = str(item).strip()
        if not text:
            continue

        run_bullet = p.add_run()
        run_bullet.text = "▸  "
        run_bullet.font.size = Pt(text_size)
        run_bullet.font.color.rgb = _rgb(*bullet_color)
        run_bullet.font.name = FONT_BODY

        run_text = p.add_run()
        run_text.text = text
        run_text.font.size = Pt(text_size)
        run_text.font.color.rgb = _rgb(*text_color)
        run_text.font.name = FONT_BODY

        p.space_after = Pt(spacing)
        p.space_before = Pt(2)
        p.line_spacing = Pt(text_size * 1.5)

    return txBox


# ─── Card Shape Helper ───────────────────────────────────────────

def _add_card(slide, left, top, width, height, fill_rgb=BG_CARD, border_rgb=None, corner_radius=Inches(0.15)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(*fill_rgb)

    if border_rgb:
        shape.line.color.rgb = _rgb(*border_rgb)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()

    shape.shadow.inherit = False

    sp_pr = shape._element.find(qn('p:spPr'))
    if sp_pr is None:
        for el in shape._element.iter():
            if el.tag.endswith('spPr'):
                sp_pr = el
                break
    if sp_pr is not None:
        prst_geom = sp_pr.find(qn('a:prstGeom'))
        if prst_geom is not None:
            av_lst = prst_geom.find(qn('a:avLst'))
            if av_lst is None:
                av_lst = etree.SubElement(prst_geom, qn('a:avLst'))
            for gd in av_lst.findall(qn('a:gd')):
                av_lst.remove(gd)
            gd = etree.SubElement(av_lst, qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', 'val 8000')

    return shape


# ─── Slide Builders ──────────────────────────────────────────────

def _build_title_slide(prs, slide_data):
    slide = prs.slides.add_slide(prs.slide_layouts[_blank_layout_idx])
    if not _use_template_bg:
        _set_three_stop_gradient_bg(slide, "06021A", "0E0830", "1A0E4A", angle=10800000)

    _add_glow_circle(slide, Inches(8), Inches(-1.5), Inches(6), ACCENT, 3000)
    _add_glow_circle(slide, Inches(-2), Inches(4), Inches(5), (124, 58, 237), 2500)
    _add_glow_circle(slide, Inches(10), Inches(5), Inches(4), TEAL, 1500)

    _add_top_accent_bar(slide)

    title_text = slide_data.get("title", "Untitled Project")
    content = slide_data.get("content", [])

    tag_box = slide.shapes.add_textbox(Inches(3.5), Inches(1.5), Inches(6.3), Inches(0.4))
    tf_tag = tag_box.text_frame
    p_tag = tf_tag.paragraphs[0]
    p_tag.alignment = PP_ALIGN.CENTER
    run_tag = p_tag.add_run()
    run_tag.text = "AI SOFTWARE COMPANY"
    run_tag.font.size = Pt(11)
    run_tag.font.bold = True
    run_tag.font.color.rgb = _rgb(*ACCENT_LIGHT)
    run_tag.font.name = FONT_HEADING
    run_tag.font.spacing = Pt(4)

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.3), Inches(2.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = str(title_text)
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = _rgb(*WHITE)
    p.font.name = FONT_HEADING
    p.line_spacing = Pt(56)

    if content:
        sub_box = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(9.3), Inches(1.5))
        tf2 = sub_box.text_frame
        tf2.word_wrap = True
        lines = content if isinstance(content, list) else [str(content)]
        for i, line in enumerate(lines):
            p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            p2.text = str(line)
            p2.font.size = Pt(18)
            p2.font.color.rgb = _rgb(*TEXT_SECONDARY)
            p2.font.name = FONT_BODY
            p2.alignment = PP_ALIGN.CENTER
            p2.space_after = Pt(6)

    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(6.2), Inches(3), Inches(0.03))
    divider.fill.solid()
    divider.fill.fore_color.rgb = _rgb(*ACCENT)
    divider.line.fill.background()
    divider.shadow.inherit = False

    if slide_data.get("speaker_notes"):
        slide.notes_slide.notes_text_frame.text = str(slide_data["speaker_notes"])

    return slide


def _build_section_divider(prs, slide_data, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[_blank_layout_idx])
    if not _use_template_bg:
        _set_gradient_bg(slide, "0A0625", "14103A", angle=8100000)

    _add_glow_circle(slide, Inches(6), Inches(0), Inches(8), ACCENT, 2000)
    _add_top_accent_bar(slide)

    title = slide_data.get("title", "")

    number_box = slide.shapes.add_textbox(MARGIN_L, Inches(2.0), Inches(2), Inches(1))
    tf_num = number_box.text_frame
    p_num = tf_num.paragraphs[0]
    run = p_num.add_run()
    run.text = f"{num:02d}"
    run.font.size = Pt(72)
    run.font.bold = True
    run.font.color.rgb = _rgb(*ACCENT)
    run.font.name = FONT_MONO

    _add_title_text(slide, title, left=MARGIN_L, top=Inches(3.2), size=42, color=WHITE)

    content = slide_data.get("content", [])
    if content:
        lines = content if isinstance(content, list) else [str(content)]
        first_line = str(lines[0]) if lines else ""
        _add_subtitle_text(slide, first_line, top=Inches(4.5), size=18, color=TEXT_SECONDARY)

    _add_bottom_line(slide)
    _add_slide_number(slide, num, total)
    _add_footer_branding(slide)

    if slide_data.get("speaker_notes"):
        slide.notes_slide.notes_text_frame.text = str(slide_data["speaker_notes"])

    return slide


def _build_content_slide(prs, slide_data, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[_blank_layout_idx])
    if not _use_template_bg:
        _set_gradient_bg(slide, "0A0625", "100C30", angle=5400000)

    _add_top_accent_bar(slide)
    _add_corner_decoration(slide, "top-right")

    title = slide_data.get("title", "")
    content = slide_data.get("content", [])

    _add_title_text(slide, title, size=32)

    title_underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_L, Inches(1.55), Inches(2.5), Inches(0.03)
    )
    title_underline.fill.solid()
    title_underline.fill.fore_color.rgb = _rgb(*ACCENT)
    title_underline.line.fill.background()
    title_underline.shadow.inherit = False

    _add_left_accent_stripe(slide, top=Inches(1.9), height=Inches(4.5))
    _add_bullet_content(slide, content)

    _add_bottom_line(slide)
    _add_slide_number(slide, num, total)
    _add_footer_branding(slide)

    if slide_data.get("speaker_notes"):
        slide.notes_slide.notes_text_frame.text = str(slide_data["speaker_notes"])

    return slide


def _build_two_column_slide(prs, slide_data, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[_blank_layout_idx])
    if not _use_template_bg:
        _set_gradient_bg(slide, "0A0625", "100C30", angle=5400000)

    _add_top_accent_bar(slide)

    title = slide_data.get("title", "")
    content = slide_data.get("content", [])

    _add_title_text(slide, title, size=32)

    title_underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_L, Inches(1.55), Inches(2.5), Inches(0.03)
    )
    title_underline.fill.solid()
    title_underline.fill.fore_color.rgb = _rgb(*ACCENT)
    title_underline.line.fill.background()
    title_underline.shadow.inherit = False

    items = content if isinstance(content, list) else [str(content)]
    mid = max(1, len(items) // 2)
    left_items = items[:mid]
    right_items = items[mid:]

    col_w = Inches(5.2)
    col_gap = Inches(0.5)
    col_top = Inches(2.0)
    col_h = Inches(4.5)

    _add_card(slide, Inches(0.7), col_top - Inches(0.15), col_w + Inches(0.3), col_h + Inches(0.3),
              fill_rgb=(18, 14, 45), border_rgb=ACCENT_DIM)
    _add_bullet_content(slide, left_items, left=Inches(1.0), top=col_top,
                        width=col_w - Inches(0.2), height=col_h, text_size=15, spacing=12)

    right_x = Inches(0.7) + col_w + col_gap
    _add_card(slide, right_x, col_top - Inches(0.15), col_w + Inches(0.3), col_h + Inches(0.3),
              fill_rgb=(18, 14, 45), border_rgb=ACCENT_DIM)
    _add_bullet_content(slide, right_items, left=right_x + Inches(0.3), top=col_top,
                        width=col_w - Inches(0.2), height=col_h,
                        bullet_color=TEAL, text_size=15, spacing=12)

    _add_bottom_line(slide)
    _add_slide_number(slide, num, total)
    _add_footer_branding(slide)

    if slide_data.get("speaker_notes"):
        slide.notes_slide.notes_text_frame.text = str(slide_data["speaker_notes"])

    return slide


def _build_stats_slide(prs, slide_data, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[_blank_layout_idx])
    if not _use_template_bg:
        _set_three_stop_gradient_bg(slide, "06021A", "0E0830", "160E40", angle=5400000)

    _add_top_accent_bar(slide)
    _add_glow_circle(slide, Inches(5), Inches(1), Inches(6), ACCENT, 1500)

    title = slide_data.get("title", "")
    content = slide_data.get("content", [])

    _add_title_text(slide, title, size=32)

    items = content if isinstance(content, list) else [str(content)]
    stat_colors = [ACCENT, TEAL, CORAL, AMBER, CYAN, ACCENT_LIGHT]
    card_w = Inches(3.4)
    card_h = Inches(2.2)
    cards_per_row = 3
    gap = Inches(0.4)
    start_x = Inches(0.9)
    start_y = Inches(2.0)

    for i, item in enumerate(items[:6]):
        row = i // cards_per_row
        col = i % cards_per_row
        x = start_x + col * (card_w + gap)
        y = start_y + row * (card_h + gap)

        _add_card(slide, x, y, card_w, card_h, fill_rgb=(18, 14, 48), border_rgb=stat_colors[i % len(stat_colors)])

        text = str(item)
        parts = text.split(":", 1) if ":" in text else [text, ""]
        label = parts[0].strip()
        value = parts[1].strip() if len(parts) > 1 else ""

        if value:
            val_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.4), card_w - Inches(0.6), Inches(0.8))
            tf_v = val_box.text_frame
            tf_v.word_wrap = True
            p_v = tf_v.paragraphs[0]
            p_v.text = value
            p_v.font.size = Pt(14)
            p_v.font.color.rgb = _rgb(*TEXT_SECONDARY)
            p_v.font.name = FONT_BODY

            lbl_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(1.3), card_w - Inches(0.6), Inches(0.5))
            tf_l = lbl_box.text_frame
            p_l = tf_l.paragraphs[0]
            p_l.text = label
            p_l.font.size = Pt(12)
            p_l.font.bold = True
            p_l.font.color.rgb = _rgb(*stat_colors[i % len(stat_colors)])
            p_l.font.name = FONT_HEADING
        else:
            val_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.35), card_w - Inches(0.6), Inches(1.5))
            tf_v = val_box.text_frame
            tf_v.word_wrap = True
            tf_v.vertical_anchor = MSO_ANCHOR.MIDDLE
            p_v = tf_v.paragraphs[0]
            p_v.text = label
            p_v.font.size = Pt(14)
            p_v.font.color.rgb = _rgb(*TEXT_SECONDARY)
            p_v.font.name = FONT_BODY

        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.2), Inches(0.08), Inches(0.08))
        dot.fill.solid()
        dot.fill.fore_color.rgb = _rgb(*stat_colors[i % len(stat_colors)])
        dot.line.fill.background()
        dot.shadow.inherit = False

    _add_bottom_line(slide)
    _add_slide_number(slide, num, total)
    _add_footer_branding(slide)

    if slide_data.get("speaker_notes"):
        slide.notes_slide.notes_text_frame.text = str(slide_data["speaker_notes"])

    return slide


def _build_comparison_slide(prs, slide_data, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[_blank_layout_idx])
    if not _use_template_bg:
        _set_gradient_bg(slide, "0A0625", "100C30", angle=5400000)

    _add_top_accent_bar(slide)

    title = slide_data.get("title", "")
    content = slide_data.get("content", [])

    _add_title_text(slide, title, size=32)

    title_underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_L, Inches(1.55), Inches(2.5), Inches(0.03)
    )
    title_underline.fill.solid()
    title_underline.fill.fore_color.rgb = _rgb(*ACCENT)
    title_underline.line.fill.background()
    title_underline.shadow.inherit = False

    items = content if isinstance(content, list) else [str(content)]

    left_w = Inches(5.5)
    right_w = Inches(5.5)
    col_top = Inches(2.0)

    left_card = _add_card(slide, Inches(0.7), col_top, left_w, Inches(4.6),
                          fill_rgb=(20, 12, 40), border_rgb=CORAL)

    lbl1 = slide.shapes.add_textbox(Inches(1.0), col_top + Inches(0.2), Inches(3), Inches(0.4))
    tf1 = lbl1.text_frame
    p1 = tf1.paragraphs[0]
    run1 = p1.add_run()
    run1.text = "EXISTING SOLUTIONS"
    run1.font.size = Pt(11)
    run1.font.bold = True
    run1.font.color.rgb = _rgb(*CORAL)
    run1.font.name = FONT_HEADING
    run1.font.spacing = Pt(2)

    mid = max(1, len(items) // 2)
    _add_bullet_content(slide, items[:mid], left=Inches(1.0), top=col_top + Inches(0.7),
                        width=left_w - Inches(0.6), height=Inches(3.5),
                        bullet_color=CORAL, text_size=15, spacing=10)

    right_card = _add_card(slide, Inches(6.5), col_top, right_w, Inches(4.6),
                           fill_rgb=(12, 20, 40), border_rgb=TEAL)

    lbl2 = slide.shapes.add_textbox(Inches(6.8), col_top + Inches(0.2), Inches(3), Inches(0.4))
    tf2 = lbl2.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "OUR APPROACH"
    run2.font.size = Pt(11)
    run2.font.bold = True
    run2.font.color.rgb = _rgb(*TEAL)
    run2.font.name = FONT_HEADING
    run2.font.spacing = Pt(2)

    _add_bullet_content(slide, items[mid:], left=Inches(6.8), top=col_top + Inches(0.7),
                        width=right_w - Inches(0.6), height=Inches(3.5),
                        bullet_color=TEAL, text_size=15, spacing=10)

    vs_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.85), Inches(3.8), Inches(0.6), Inches(0.6))
    vs_circle.fill.solid()
    vs_circle.fill.fore_color.rgb = _rgb(*BG_DEEP)
    vs_circle.line.color.rgb = _rgb(*ACCENT)
    vs_circle.line.width = Pt(2)
    vs_circle.shadow.inherit = False

    vs_box = slide.shapes.add_textbox(Inches(5.85), Inches(3.88), Inches(0.6), Inches(0.45))
    tf_vs = vs_box.text_frame
    tf_vs.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_vs = tf_vs.paragraphs[0]
    p_vs.alignment = PP_ALIGN.CENTER
    run_vs = p_vs.add_run()
    run_vs.text = "VS"
    run_vs.font.size = Pt(11)
    run_vs.font.bold = True
    run_vs.font.color.rgb = _rgb(*ACCENT)
    run_vs.font.name = FONT_HEADING

    _add_bottom_line(slide)
    _add_slide_number(slide, num, total)
    _add_footer_branding(slide)

    if slide_data.get("speaker_notes"):
        slide.notes_slide.notes_text_frame.text = str(slide_data["speaker_notes"])

    return slide


def _build_product_slide(prs, slide_data, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[_blank_layout_idx])
    if not _use_template_bg:
        _set_three_stop_gradient_bg(slide, "06021A", "0E0830", "160E40", angle=8100000)

    _add_top_accent_bar(slide)
    _add_glow_circle(slide, Inches(9), Inches(-1), Inches(5), ACCENT, 2000)

    title = slide_data.get("title", "")
    content = slide_data.get("content", [])

    _add_title_text(slide, title, size=32)

    title_underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_L, Inches(1.55), Inches(2.5), Inches(0.03)
    )
    title_underline.fill.solid()
    title_underline.fill.fore_color.rgb = _rgb(*ACCENT)
    title_underline.line.fill.background()
    title_underline.shadow.inherit = False

    items = content if isinstance(content, list) else [str(content)]
    feature_icons = ["◆", "◈", "◇", "▣", "◉", "▹"]
    feature_colors = [ACCENT, TEAL, AMBER, CYAN, CORAL, ACCENT_LIGHT]

    card_w = Inches(5.4)
    card_h = Inches(1.2)
    start_y = Inches(1.9)
    gap = Inches(0.15)

    for i, item in enumerate(items[:6]):
        col = i % 2
        row = i // 2
        x = Inches(0.7) + col * (card_w + Inches(0.3))
        y = start_y + row * (card_h + gap)
        color = feature_colors[i % len(feature_colors)]
        icon = feature_icons[i % len(feature_icons)]

        _add_card(slide, x, y, card_w, card_h, fill_rgb=(16, 12, 42), border_rgb=color)

        icon_box = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.25), Inches(0.4), Inches(0.4))
        tf_icon = icon_box.text_frame
        p_icon = tf_icon.paragraphs[0]
        p_icon.text = icon
        p_icon.font.size = Pt(16)
        p_icon.font.color.rgb = _rgb(*color)

        text_box = slide.shapes.add_textbox(x + Inches(0.65), y + Inches(0.2), card_w - Inches(1.0), card_h - Inches(0.3))
        tf_text = text_box.text_frame
        tf_text.word_wrap = True
        tf_text.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_text = tf_text.paragraphs[0]
        p_text.text = str(item)
        p_text.font.size = Pt(14)
        p_text.font.color.rgb = _rgb(*TEXT_SECONDARY)
        p_text.font.name = FONT_BODY

    _add_bottom_line(slide)
    _add_slide_number(slide, num, total)
    _add_footer_branding(slide)

    if slide_data.get("speaker_notes"):
        slide.notes_slide.notes_text_frame.text = str(slide_data["speaker_notes"])

    return slide


def _build_thank_you_slide(prs, slide_data, num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[_blank_layout_idx])
    if not _use_template_bg:
        _set_three_stop_gradient_bg(slide, "06021A", "0E0830", "1A0E4A", angle=10800000)

    _add_glow_circle(slide, Inches(4), Inches(0), Inches(8), ACCENT, 2500)
    _add_glow_circle(slide, Inches(-1), Inches(3), Inches(5), (124, 58, 237), 2000)
    _add_glow_circle(slide, Inches(9), Inches(4), Inches(4), TEAL, 1500)

    _add_top_accent_bar(slide)

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11.3), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Thank You"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = _rgb(*ACCENT)
    p.font.name = FONT_HEADING

    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.6), Inches(2.3), Inches(0.03)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = _rgb(*ACCENT_DIM)
    divider.line.fill.background()
    divider.shadow.inherit = False

    content = slide_data.get("content", ["Questions?"])
    lines = content if isinstance(content, list) else [str(content)]

    sub_box = slide.shapes.add_textbox(Inches(2), Inches(4.0), Inches(9.3), Inches(2))
    tf2 = sub_box.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(lines):
        p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p2.text = str(line)
        p2.font.size = Pt(18)
        p2.font.color.rgb = _rgb(*TEXT_SECONDARY)
        p2.font.name = FONT_BODY
        p2.alignment = PP_ALIGN.CENTER
        p2.space_after = Pt(8)

    brand_box = slide.shapes.add_textbox(Inches(3.5), Inches(6.2), Inches(6.3), Inches(0.4))
    tf3 = brand_box.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "BUILT WITH AI SOFTWARE COMPANY"
    run3.font.size = Pt(9)
    run3.font.bold = True
    run3.font.color.rgb = _rgb(*TEXT_MUTED)
    run3.font.name = FONT_HEADING
    run3.font.spacing = Pt(4)

    _add_slide_number(slide, num, total)

    if slide_data.get("speaker_notes"):
        slide.notes_slide.notes_text_frame.text = str(slide_data["speaker_notes"])

    return slide


# ─── Slide Type Router ───────────────────────────────────────────

SLIDE_LAYOUT_MAP = {
    1: "title",
    2: "content",
    3: "content",
    4: "section_divider",
    5: "content",
    6: "product",
    7: "stats",
    8: "comparison",
    9: "two_column",
    10: "thank_you",
}

TITLE_KEYWORD_MAP = {
    "introduction": "content",
    "problem": "content",
    "insight": "section_divider",
    "solution": "content",
    "product": "product",
    "overview": "product",
    "viability": "stats",
    "impact": "stats",
    "sustainability": "stats",
    "competitive": "comparison",
    "prior art": "comparison",
    "analysis": "comparison",
    "marketing": "two_column",
    "branding": "two_column",
    "strategy": "two_column",
    "thank": "thank_you",
}


def _detect_slide_type(slide_data, slide_num, total_slides):
    title = slide_data.get("title", "").lower()

    if slide_num == 1:
        return "title"
    if slide_num == total_slides and "thank" in title:
        return "thank_you"

    for keyword, layout in TITLE_KEYWORD_MAP.items():
        if keyword in title:
            return layout

    return SLIDE_LAYOUT_MAP.get(slide_num, "content")


BUILDERS = {
    "title": _build_title_slide,
    "content": _build_content_slide,
    "section_divider": _build_section_divider,
    "two_column": _build_two_column_slide,
    "stats": _build_stats_slide,
    "comparison": _build_comparison_slide,
    "product": _build_product_slide,
    "thank_you": _build_thank_you_slide,
}


# ─── Main Generator ─────────────────────────────────────────────

def generate_pptx(project_id: str, ppt_output: dict) -> str:
    PROJECTS_DIR.mkdir(exist_ok=True)

    prs = _create_presentation(project_id)

    slides_data = ppt_output.get("slides", [])
    pitch = ppt_output.get("pitch", {})

    if not slides_data and pitch:
        slides_data = _pitch_to_slides(pitch)

    if not slides_data:
        slides_data = [{"title": "Untitled", "content": ["No slide data generated"], "speaker_notes": ""}]

    total = len(slides_data)

    for i, slide_data in enumerate(slides_data):
        num = i + 1
        slide_type = _detect_slide_type(slide_data, num, total)
        builder = BUILDERS.get(slide_type, _build_content_slide)

        if slide_type == "title":
            builder(prs, slide_data)
        else:
            builder(prs, slide_data, num, total)

    output_path = PROJECTS_DIR / f"{project_id}_presentation.pptx"
    prs.save(str(output_path))
    return str(output_path)


def _pitch_to_slides(pitch: dict) -> list[dict]:
    slides = []

    hook = pitch.get("hook", "Our Idea")
    slides.append({
        "title": str(hook),
        "content": ["AI Software Company"],
        "speaker_notes": "Opening — introduce the team and idea."
    })

    problem = pitch.get("problem", "")
    slides.append({
        "title": "Introduction",
        "content": [str(problem)] if problem else ["Introducing our solution"],
        "speaker_notes": "Brief context on what we're building."
    })

    slides.append({
        "title": "Problem Statement",
        "content": [str(problem)] if problem else ["The core challenge"],
        "speaker_notes": "Explain the pain point clearly."
    })

    solution = pitch.get("solution", "")
    slides.append({
        "title": "Our Insight",
        "content": [str(solution)] if solution else ["What we discovered"],
        "speaker_notes": "The unique observation that drives our approach."
    })

    slides.append({
        "title": "Solution Overview",
        "content": [str(solution)] if solution else ["How we solve it"],
        "speaker_notes": "High-level solution."
    })

    features = pitch.get("key_features", [])
    tech = pitch.get("tech_highlights", "")
    product_content = []
    if features:
        product_content.extend([str(f) for f in features])
    if tech:
        product_content.append(str(tech) if isinstance(tech, str) else ", ".join(str(t) for t in tech))
    slides.append({
        "title": "Product Overview",
        "content": product_content or ["Key features and tech stack"],
        "speaker_notes": "Walk through features."
    })

    slides.append({
        "title": "Viability, Sustainability & Impact",
        "content": ["Market visibility", "Sustainability plan", "Long-term survival scope", "Impact on target users"],
        "speaker_notes": "Why this idea can survive and grow."
    })

    slides.append({
        "title": "Competitive Analysis & Prior Art",
        "content": ["Existing solutions comparison", "Our differentiation"],
        "speaker_notes": "How we stand out."
    })

    future = pitch.get("future_vision", "")
    slides.append({
        "title": "Marketing & Branding Strategy",
        "content": [str(future)] if future else ["Go-to-market plan", "Growth strategy"],
        "speaker_notes": "How we reach users."
    })

    slides.append({
        "title": "Thank You",
        "content": ["Questions?"],
        "speaker_notes": "Open the floor."
    })

    return slides


def get_pptx_path(project_id: str) -> str | None:
    path = PROJECTS_DIR / f"{project_id}_presentation.pptx"
    if path.exists():
        return str(path)
    return None
